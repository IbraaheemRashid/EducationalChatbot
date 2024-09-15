import os
import re
import nltk
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

nltk.download('punkt_tab')

def parse_pdf(file_path):
    try:
        text = ""
        with open(file_path, "rb") as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
        return text
    except Exception as e:
        print(f"Error parsing PDF file: {e}")
        return None

def parse_pptx(file_path):
    try:
        text = ""
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error parsing PPTX file: {e}")
        return None
    
def parse_docx(file_path):
    try:
        text = ""
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        print(f"Error parsing DOCX file: {e}")
        return None

def clean_text(text):
    try:
        text = text.strip()
        text = re.sub(r'\s+', ' ', text)
        text = text.encode("ascii", "ignore").decode()
        text = re.sub(r'[^\w\s.,!?;]', '', text)
        return text
    except Exception as e:
        print(f"Error cleaning text: {e}")
        return None

def chunk_text(text, max_words=150):
    try:
        sentences = nltk.sent_tokenize(text)
        chunks = []
        current_chunk = []
        current_length = 0

        for sentence in sentences:
            words = sentence.split()
            if current_length + len(words) <= max_words:
                current_chunk.append(sentence)
                current_length += len(words)
            else:
                chunks.append(" ".join(current_chunk))
                current_chunk = [sentence]
                current_length = len(words)

        if current_chunk:
            chunks.append(" ".join(current_chunk))

        return chunks
    except Exception as e:
        print(f"Error chunking text: {e}")
        return []

def parse_document(file_path):
    import os
    extension = os.path.splitext(file_path)[1].lower()
    print(f"Parsing file with extension: {extension}")
    print(f"File path provided: {file_path}")

    if extension == '.pdf':
        return parse_pdf(file_path)
    elif extension == '.pptx':
        return parse_pptx(file_path)
    elif extension == '.docx':
        return parse_docx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {extension}")

