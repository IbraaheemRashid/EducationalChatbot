import os
import re
import nltk
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from nltk.tokenize import sent_tokenize

nltk.download('punkt', quiet=True)

def parse_pdf(file_path):
    try:
        text = ""
        with open(file_path, "rb") as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
        return text
    except Exception as e:
        print(f"Error parsing PDF file: {e}")
        return None

def parse_pptx(file_path):
    try:
        text = ""
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error parsing PPTX file: {e}")
        return None
    
def parse_docx(file_path):
    try:
        text = ""
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        print(f"Error parsing DOCX file: {e}")
        return None

def clean_text(text):
    try:
        text = text.strip()
        text = re.sub(r'\s+', ' ', text)
        text = text.encode("ascii", "ignore").decode()
        text = re.sub(r'[^\w\s.,!?;]', '', text)
        return text
    except Exception as e:
        print(f"Error cleaning text: {e}")
        return None

def chunk_text(text, chunk_size=100, overlap=20):
    words = text.split()
    chunks = []
    start = 0
    
    while start < len(words):
        end = start + chunk_size
        chunk = ' '.join(words[start:end])
        if chunk not in chunks:
            chunks.append(chunk)
        start += chunk_size - overlap
    
    return chunks


def parse_document(file_path):
    extension = os.path.splitext(file_path)[1].lower()

    if extension == '.pdf':
        return parse_pdf(file_path)
    elif extension == '.pptx':
        return parse_pptx(file_path)
    elif extension == '.docx':
        return parse_docx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {extension}")

def print_chunk_debug_info(chunks):
    print(f"Number of chunks: {len(chunks)}")
    print("First few chunks:")
    for i, chunk in enumerate(chunks[:3]):
        print(f"Chunk {i}: {chunk[:100]}...")  # Print first 100 characters of each chunk